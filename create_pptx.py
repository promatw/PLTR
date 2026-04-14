from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette (Ocean / Midnight) ──────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # very dark navy
NAVY       = RGBColor(0x06, 0x5A, 0x82)   # deep blue
TEAL       = RGBColor(0x1C, 0x72, 0x93)   # teal
MINT       = RGBColor(0x02, 0xC3, 0x9A)   # mint accent
LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF8)   # off-white
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
GOLD       = RGBColor(0xF5, 0xC5, 0x18)
RED_WARN   = RGBColor(0xE5, 0x3E, 0x3E)
CARD_BG    = RGBColor(0x0F, 0x2B, 0x40)   # dark card

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font_face="Calibri", wrap=True, margin=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if margin is not None:
        tf.margin_left = Inches(margin)
        tf.margin_right = Inches(margin)
        tf.margin_top = Inches(margin)
        tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_face
    return txBox

def add_para(tf, text, font_size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font_face="Calibri", space_before=None):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_face
    return p

def add_bullet_box(slide, items, x, y, w, h, font_size=13, color=DARK_TEXT,
                   font_face="Calibri", bullet_color=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_face
    return txBox

def add_card(slide, x, y, w, h, title, body_lines, title_color=WHITE,
             body_color=None, bg_color=CARD_BG, accent_color=MINT):
    """Adds a dark card with accent left border, title, and bullet body."""
    if body_color is None:
        body_color = RGBColor(0xCB, 0xD5, 0xE1)
    # Card background
    add_rect(slide, x, y, w, h, bg_color)
    # Accent border
    add_rect(slide, x, y, 0.07, h, accent_color)
    # Title
    add_text(slide, title, x+0.15, y+0.08, w-0.2, 0.35,
             13, title_color, bold=True, font_face="Calibri")
    # Body
    body_box = slide.shapes.add_textbox(
        Inches(x+0.15), Inches(y+0.42), Inches(w-0.25), Inches(h-0.55))
    body_box.word_wrap = True
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.color.rgb = body_color
        run.font.name = "Calibri"

def stat_block(slide, x, y, w, h, number, label, num_color=MINT, bg_color=CARD_BG):
    add_rect(slide, x, y, w, h, bg_color)
    add_rect(slide, x, y, w, 0.05, num_color)  # top accent line
    add_text(slide, number, x, y+0.12, w, 0.55,
             28, num_color, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")
    add_text(slide, label, x, y+0.68, w, 0.55,
             10, RGBColor(0xCB, 0xD5, 0xE1), bold=False,
             align=PP_ALIGN.CENTER, font_face="Calibri")

# ── Presentation Setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # completely blank

SLIDE_W = 13.33
SLIDE_H = 7.5

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK_BG)

# Bottom teal bar
add_rect(s1, 0, 6.9, SLIDE_W, 0.6, NAVY)
# Mint accent stripe (left)
add_rect(s1, 0, 0, 0.35, SLIDE_H, MINT)
# Teal diagonal accent block (right decorative)
add_rect(s1, 10.5, 0, 2.83, SLIDE_H, CARD_BG)
add_rect(s1, 10.5, 0, 0.06, SLIDE_H, TEAL)

# Category tag
add_rect(s1, 0.7, 0.8, 2.8, 0.42, MINT)
add_text(s1, "投資分析  ·  AI產業", 0.75, 0.83, 2.7, 0.36,
         11, DARK_BG, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")

# Main title
add_text(s1, "暴跌 12%！", 0.7, 1.4, 9.5, 1.0,
         44, GOLD, bold=True, font_face="Calibri")
add_text(s1, "Palantir 的生存危機", 0.7, 2.3, 9.5, 0.8,
         36, WHITE, bold=True, font_face="Calibri")

# Subtitle
add_text(s1,
    "大空頭 Michael Burry 做空邏輯 × Anthropic 的降維打擊\n"
    "蒂爾死守的最後一張牌，正被供應商「生吞」？",
    0.7, 3.2, 9.3, 1.1,
    16, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")

# Divider line
add_rect(s1, 0.7, 4.45, 4.5, 0.05, TEAL)

# Bottom tags
add_text(s1, "PLTR  ·  PLTR vs Anthropic  ·  AI 估值攻防", 0.7, 4.6, 7, 0.4,
         12, MUTED, font_face="Calibri")

# Right panel content
add_text(s1, "[ 統計 ]", 11.0, 1.5, 1.8, 1.5, 60, TEAL, align=PP_ALIGN.CENTER, font_face="Calibri")
add_text(s1, "影片內容重點整理", 10.55, 3.1, 2.6, 0.4,
         11, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, font_face="Calibri")
add_text(s1, "繁體中文版", 10.55, 3.5, 2.6, 0.35,
         13, MINT, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")

# Footer
add_text(s1, "資料來源：YouTube 影片分析  |  NotebookLM 整理  |  2026.04",
         0.5, 7.0, 10, 0.35, 10, MUTED, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 事件背景
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, LIGHT_BG)

# Header bar
add_rect(s2, 0, 0, SLIDE_W, 1.1, NAVY)
add_rect(s2, 0, 1.1, SLIDE_W, 0.06, MINT)
add_text(s2, "01  /  事件背景", 0.5, 0.1, 3, 0.35, 11, MINT, bold=True, font_face="Calibri")
add_text(s2, "為何 PLTR 逆勢暴跌？", 0.5, 0.38, 9, 0.6, 28, WHITE, bold=True, font_face="Calibri")

# Left block: context
add_rect(s2, 0.4, 1.4, 5.8, 5.5, WHITE)
add_rect(s2, 0.4, 1.4, 0.08, 5.5, RED_WARN)

add_text(s2, "市場背景", 0.65, 1.55, 5.3, 0.4, 14, DARK_TEXT, bold=True, font_face="Calibri")

ctx_lines = [
    "美股大盤近期反彈，納斯達克上漲 2%",
    "大部分科技股同步大漲",
    "Palantir 卻在兩天內逆勢暴跌超過 12%",
    "從年初高點累計下跌超過 35%",
    "",
    "這說明不是大盤拖累的——",
    "是 Palantir 自身出了問題！",
]
y_pos = 2.05
for line in ctx_lines:
    color = RED_WARN if "!" in line else DARK_TEXT
    bold = "!" in line
    add_text(s2, line, 0.65, y_pos, 5.3, 0.38, 13 if not bold else 14,
             color, bold=bold, font_face="Calibri")
    y_pos += 0.42

# Right block: two stat cards
add_text(s2, "關鍵數字", 6.7, 1.55, 5.9, 0.4, 14, DARK_TEXT, bold=True, font_face="Calibri")

stat_block(s2, 6.7, 2.05, 2.7, 1.55, "−12%", "兩天累計跌幅", num_color=RED_WARN)
stat_block(s2, 9.6, 2.05, 2.7, 1.55, "−35%", "距年初高點跌幅", num_color=RED_WARN)

stat_block(s2, 6.7, 3.75, 2.7, 1.55, "$3,000億+", "市值（高峰期）", num_color=MINT)
stat_block(s2, 9.6, 3.75, 2.7, 1.55, "70x", "市銷率 P/S\n（標普500最貴之一）", num_color=GOLD)

# Quote
add_rect(s2, 6.7, 5.45, 5.6, 1.3, CARD_BG)
add_rect(s2, 6.7, 5.45, 0.07, 1.3, GOLD)
add_text(s2, '"Anthropic 正在吃 Palantir 的午餐"', 6.85, 5.55, 5.3, 0.45,
         14, GOLD, bold=True, italic=True, font_face="Calibri")
add_text(s2, "— Michael Burry，大空頭，做空次貸危機獲利 7.25 億美元的傳奇投資人",
         6.85, 6.05, 5.3, 0.55, 11, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")

# Footer
add_rect(s2, 0, 7.15, SLIDE_W, 0.35, NAVY)
add_text(s2, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s2, "2 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 三大觸發事件
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, LIGHT_BG)

add_rect(s3, 0, 0, SLIDE_W, 1.1, DARK_BG)
add_rect(s3, 0, 1.1, SLIDE_W, 0.06, MINT)
add_text(s3, "02  /  觸發因素", 0.5, 0.1, 3, 0.35, 11, MINT, bold=True, font_face="Calibri")
add_text(s3, "三大關鍵事件同時爆發", 0.5, 0.38, 9, 0.6, 28, WHITE, bold=True, font_face="Calibri")

# Three event cards
events = [
    ("事件 ① ", "Michael Burry 公開喊話",
     ["週二在 X 發文：「Anthropic 正在吃 Palantir 午餐」",
      "持有大量 2027年到期、行使價 $50 的看跌期權",
      "從 2024年11月起就開始建立做空倉位",
      "引用 RAMP AI Index 最新企業消費數據"],
     GOLD),
    ("事件 ② ", "Anthropic 發布 Managed Agents",
     ["企業不再需要搭建複雜 AI 平台",
      "只需描述任務，AI 助手直接完成",
      "按任務收費，一個 AI 助手抵 100 人工作",
      "直接繞過 Palantir 提供的「中間平台層」"],
     MINT),
    ("事件 ③ ", "Anthropic 發布 Claude Methods 模型",
     ["強大到不敢公開發布的網路安全 AI 模型",
      "找到藏了 27年的 OpenBSD 安全漏洞",
      "僅提供 12家核心夥伴（Apple、Microsoft…）",
      "Palantir 軍方系統仍在使用舊版 Claude"],
     RED_WARN),
]

for i, (num, title, bullets, accent) in enumerate(events):
    x = 0.4 + i * 4.3
    w = 4.1
    add_rect(s3, x, 1.35, w, 5.75, WHITE)
    add_rect(s3, x, 1.35, w, 0.08, accent)

    # Number tag
    add_rect(s3, x+0.15, 1.6, 0.8, 0.38, accent)
    add_text(s3, num, x+0.15, 1.6, 0.8, 0.38, 11, DARK_BG, bold=True,
             align=PP_ALIGN.CENTER, font_face="Calibri")

    add_text(s3, title, x+1.05, 1.63, w-1.2, 0.38, 13, DARK_TEXT,
             bold=True, font_face="Calibri")

    add_rect(s3, x+0.15, 2.12, w-0.3, 0.03, RGBColor(0xE2, 0xE8, 0xF0))

    y_b = 2.28
    for bullet in bullets:
        add_text(s3, f"▸  {bullet}", x+0.2, y_b, w-0.35, 0.55,
                 12, DARK_TEXT, font_face="Calibri")
        y_b += 0.58

# Impact label
add_rect(s3, 0.4, 6.9, SLIDE_W - 0.8, 0.45, DARK_BG)
add_text(s3, "[ 核心 ]  三個事件加在一起，直接撼動了 Palantir 的核心商業邏輯",
         0.6, 6.93, 12, 0.38, 13, GOLD, bold=True, font_face="Calibri")

# Footer
add_rect(s3, 0, 7.15, SLIDE_W, 0.35, NAVY)
add_text(s3, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s3, "3 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Managed Agents：商業攻勢
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, DARK_BG)

add_rect(s4, 0, 0, SLIDE_W, 1.1, CARD_BG)
add_rect(s4, 0, 1.1, SLIDE_W, 0.06, TEAL)
add_text(s4, "03  /  Anthropic 商業攻勢", 0.5, 0.1, 4.5, 0.35, 11, TEAL, bold=True, font_face="Calibri")
add_text(s4, "Managed Agents：AI 從「工具」變「工人」", 0.5, 0.38, 11, 0.6, 26, WHITE, bold=True, font_face="Calibri")

# Left: analogy explanation
add_rect(s4, 0.4, 1.3, 6.2, 5.9, CARD_BG)
add_rect(s4, 0.4, 1.3, 0.07, 5.9, TEAL)

add_text(s4, "[ 比喻 ]  用「餐廳」比喻理解商業模式差異", 0.6, 1.45, 5.8, 0.45,
         13, TEAL, bold=True, font_face="Calibri")

# Palantir model
add_rect(s4, 0.6, 2.0, 5.8, 2.0, rgb(0x06, 0x2A, 0x40))
add_rect(s4, 0.6, 2.0, 5.8, 0.07, GOLD)
add_text(s4, "❌  Palantir 舊模式：企業自建廚房", 0.75, 2.1, 5.5, 0.38,
         13, GOLD, bold=True, font_face="Calibri")
old_lines = [
    "花大錢搭建 AI 平台（廚房）、培訓員工（廚師）",
    "設定菜單（工作流程），員工按人頭付費使用",
    "成本高、建置慢，往往需要數月以上",
]
y_ = 2.55
for l in old_lines:
    add_text(s4, f"  {l}", 0.75, y_, 5.5, 0.42, 12, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_ += 0.42

# Anthropic model
add_rect(s4, 0.6, 4.15, 5.8, 2.1, rgb(0x02, 0x3F, 0x35))
add_rect(s4, 0.6, 4.15, 5.8, 0.07, MINT)
add_text(s4, "✅  Anthropic 新模式：AI 廚師直接到府服務", 0.75, 4.25, 5.5, 0.38,
         13, MINT, bold=True, font_face="Calibri")
new_lines = [
    "描述需求，AI 助手自動讀文件、查數據、寫報告",
    "按任務收費，一個 AI 助手處理 100 人的工作",
    "一週內即可部署（Rakuten 實際案例）",
]
y_ = 4.7
for l in new_lines:
    add_text(s4, f"  {l}", 0.75, y_, 5.5, 0.42, 12, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_ += 0.42

# Right: impact analysis
add_text(s4, "對 Palantir 的衝擊", 7.0, 1.45, 5.9, 0.4, 14, WHITE, bold=True, font_face="Calibri")

impacts = [
    ("核心邏輯被動搖",
     "Palantir 賣的是「基礎設施」\nAnthropic 賣的是「結果」\n企業真正需要的是結果！", TEAL),
    ("按席位收費模式危機",
     "傳統軟體按人頭收費\nAI 助手按任務收費\n平台層的必要性正在消失", GOLD),
    ("實際案例已發生",
     "Rakuten：一週部署三部門 AI\nLinear：AI 直接接管專案管理\nNotion：AI 在工作流程裡協作", MINT),
]
y_i = 1.95
for title_i, body_i, accent_i in impacts:
    add_rect(s4, 7.0, y_i, 5.9, 1.6, rgb(0x0F, 0x2B, 0x40))
    add_rect(s4, 7.0, y_i, 0.07, 1.6, accent_i)
    add_text(s4, title_i, 7.15, y_i+0.1, 5.6, 0.38, 12, accent_i, bold=True, font_face="Calibri")
    add_text(s4, body_i, 7.15, y_i+0.52, 5.6, 1.0, 11,
             RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_i += 1.72

# Footer
add_rect(s4, 0, 7.15, SLIDE_W, 0.35, CARD_BG)
add_text(s4, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s4, "4 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Claude Methods：技術代差
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, DARK_BG)

add_rect(s5, 0, 0, SLIDE_W, 1.1, CARD_BG)
add_rect(s5, 0, 1.1, SLIDE_W, 0.06, RED_WARN)
add_text(s5, "04  /  Anthropic 技術代差", 0.5, 0.1, 4.5, 0.35, 11, RED_WARN, bold=True, font_face="Calibri")
add_text(s5, "Claude Methods：強大到不敢公開的模型", 0.5, 0.38, 11, 0.6, 26, WHITE, bold=True, font_face="Calibri")

# Left column: Methods capabilities
add_text(s5, "[ 安全 ]  Claude Methods 的能力", 0.4, 1.3, 6.2, 0.4, 14, RED_WARN, bold=True, font_face="Calibri")

caps = [
    ("找到隱藏 27年的安全漏洞",
     "在主要作業系統和瀏覽器中發現關鍵安全漏洞，"
     "其中最老的漏洞藏在 OpenBSD 中長達 27 年"),
    ("自動組成攻擊鏈",
     "將多個漏洞串聯成完整攻擊鏈，"
     "可讓普通用戶帳號取得整台 Linux 伺服器最高權限"),
    ("自主逃脫虛擬環境",
     "在封閉測試環境中，模型自主找到方法突破隔離，"
     "獲取網路權限並主動聯絡研究員"),
    ("能力等同頂尖安全專家",
     "Anthropic 安全負責人表示：此模型能力媲美全球"
     "僅有幾百人的頂尖網路安全專家"),
]
y_c = 1.85
for cap_title, cap_body in caps:
    add_rect(s5, 0.4, y_c, 6.2, 1.2, CARD_BG)
    add_rect(s5, 0.4, y_c, 0.07, 1.2, RED_WARN)
    add_text(s5, cap_title, 0.58, y_c+0.08, 5.9, 0.38, 12, RED_WARN, bold=True, font_face="Calibri")
    add_text(s5, cap_body, 0.58, y_c+0.48, 5.9, 0.65, 11,
             RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_c += 1.3

# Right column: Implications for Palantir
add_text(s5, "[ 衝擊 ]  對 Palantir 的雙重打擊", 6.9, 1.3, 6.0, 0.4, 14, GOLD, bold=True, font_face="Calibri")

impl = [
    ("打擊 ①：技術代差擴大",
     "Palantir 軍方 Maven 系統底層依賴 Anthropic 的模型，"
     "但目前仍使用舊版 Claude。\n"
     "而 Methods 模型比現有版本強大許多，且 Palantir 無法獲取。",
     GOLD),
    ("打擊 ②：供應商變成競爭對手",
     "Anthropic 決定不對外公開 Methods，"
     "僅提供 Apple、Microsoft 等 12 家核心夥伴。\n"
     "供應商正在繞過 Palantir 直接服務企業客戶。",
     MINT),
    ("打擊 ③：五角大廈關係惡化",
     "Anthropic 因安全限制爭議被五角大廈列為「供應鏈風險」，"
     "雙方關係緊張。\n"
     "這讓 Palantir 難以升級底層技術棧。",
     RED_WARN),
]
y_i = 1.85
for imp_title, imp_body, accent_i in impl:
    add_rect(s5, 6.9, y_i, 6.0, 1.65, CARD_BG)
    add_rect(s5, 6.9, y_i, 0.07, 1.65, accent_i)
    add_text(s5, imp_title, 7.08, y_i+0.1, 5.7, 0.38, 12, accent_i, bold=True, font_face="Calibri")
    add_text(s5, imp_body, 7.08, y_i+0.5, 5.7, 1.08, 11,
             RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_i += 1.78

# Footer
add_rect(s5, 0, 7.15, SLIDE_W, 0.35, CARD_BG)
add_text(s5, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s5, "5 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Anthropic 爆炸性成長數據
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, LIGHT_BG)

add_rect(s6, 0, 0, SLIDE_W, 1.1, NAVY)
add_rect(s6, 0, 1.1, SLIDE_W, 0.06, MINT)
add_text(s6, "05  /  Anthropic 關鍵數據", 0.5, 0.1, 4.5, 0.35, 11, MINT, bold=True, font_face="Calibri")
add_text(s6, "爆炸性成長：B2B 軟體歷史上前所未有", 0.5, 0.38, 11, 0.6, 26, WHITE, bold=True, font_face="Calibri")

# Top: 4 big stats
stats_top = [
    ("70%", "企業首次購買 AI 工具選擇 Anthropic\n（2025年前為 OpenAI 領先）"),
    ("$300億", "年化收入 ARR\n（15個月前僅 $1億）"),
    ("10倍", "使用付費 Anthropic 的企業數\n（一年內成長幅度）"),
    (">20%", "預測年底 GitHub 公開代碼\n由 Claude 撰寫比例"),
]
for i, (num, label) in enumerate(stats_top):
    x_s = 0.4 + i * 3.13
    add_rect(s6, x_s, 1.35, 2.9, 1.8, DARK_BG)
    add_rect(s6, x_s, 1.35, 2.9, 0.08, MINT)
    add_text(s6, num, x_s, 1.5, 2.9, 0.7, 30, MINT, bold=True,
             align=PP_ALIGN.CENTER, font_face="Calibri")
    add_text(s6, label, x_s+0.1, 2.2, 2.7, 0.85, 11,
             RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER, font_face="Calibri")

# Bottom: timeline + context
add_rect(s6, 0.4, 3.35, 8.0, 3.9, WHITE)
add_rect(s6, 0.4, 3.35, 0.07, 3.9, MINT)
add_text(s6, "[ 成長 ]  成長時間軸", 0.6, 3.5, 7.5, 0.4, 13, DARK_TEXT, bold=True, font_face="Calibri")

timeline = [
    ("2024年初", "$1億", "年化收入（ARR）", MUTED),
    ("2025年底", "$90億", "年化收入（僅 15 個月後）", TEAL),
    ("2026年4月", "$300億", "年化收入（再過 4 個月翻逾 3 倍）", MINT),
]
y_t = 4.05
for date, amt, label, color in timeline:
    add_rect(s6, 0.65, y_t+0.1, 0.12, 0.12, color)
    add_text(s6, date, 0.9, y_t, 1.5, 0.38, 12, MUTED, font_face="Calibri")
    add_text(s6, amt, 2.5, y_t, 1.5, 0.38, 16, color, bold=True, font_face="Calibri")
    add_text(s6, label, 4.1, y_t, 4.2, 0.38, 12, DARK_TEXT, font_face="Calibri")
    if y_t < 5.5:
        add_rect(s6, 0.7, y_t+0.42, 0.02, 0.55, RGBColor(0xCB, 0xD5, 0xE1))
    y_t += 0.9

add_text(s6, "整個軟體行業 50 年歷史裡，從未有任何一家公司在這個體量上以此速度增長。",
         0.6, 6.7, 7.6, 0.45, 12, DARK_TEXT, italic=True, font_face="Calibri")

# Right box: RAMP note + GitHub
add_rect(s6, 8.65, 3.35, 4.3, 3.9, DARK_BG)
add_rect(s6, 8.65, 3.35, 4.3, 0.07, GOLD)
add_text(s6, "RAMP AI Index 數據來源說明", 8.8, 3.5, 4.0, 0.4, 12, GOLD, bold=True, font_face="Calibri")
note_lines = [
    "RAMP 是美國企業信用卡平台，",
    "追蹤數萬家公司在 AI 上的花費",
    "",
    "⚠️ 侷限性：",
    "• 主要追蹤中小企業信用卡",
    "• 大企業 AWS 長期合約不在內",
    "• 可能高估 Anthropic 優勢",
    "",
    "GitHub 數據：",
    "目前 4% 公開代碼由 Claude 撰寫",
    "預計年底突破 20%",
]
y_n = 4.0
for nl in note_lines:
    color_n = GOLD if "⚠️" in nl else (MINT if "GitHub" in nl else RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s6, nl, 8.8, y_n, 4.0, 0.38, 11, color_n, font_face="Calibri")
    y_n += 0.33

# Footer
add_rect(s6, 0, 7.15, SLIDE_W, 0.35, NAVY)
add_text(s6, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s6, "6 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Palantir 的基本面
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_bg(s7, LIGHT_BG)

add_rect(s7, 0, 0, SLIDE_W, 1.1, DARK_BG)
add_rect(s7, 0, 1.1, SLIDE_W, 0.06, GOLD)
add_text(s7, "06  /  Palantir 基本面", 0.5, 0.1, 4.5, 0.35, 11, GOLD, bold=True, font_face="Calibri")
add_text(s7, "業務持續擴張：數字亮眼，護城河深厚", 0.5, 0.38, 11, 0.6, 26, WHITE, bold=True, font_face="Calibri")

# Top stat row: financials
fin_stats = [
    ("70%", "Q4 2025 營收同比增長\n（連續第 10 季加速）"),
    ("137%", "美國商業部門收入增長"),
    ("127", "Rule of 40 指標\n（行業健康線為 40）"),
    ("$43億", "單季合約總價值\n（創歷史新高）"),
]
for i, (num, label) in enumerate(fin_stats):
    x_s = 0.4 + i * 3.13
    add_rect(s7, x_s, 1.3, 2.9, 1.8, WHITE)
    add_rect(s7, x_s, 1.3, 2.9, 0.08, GOLD)
    add_text(s7, num, x_s, 1.45, 2.9, 0.7, 28, GOLD, bold=True,
             align=PP_ALIGN.CENTER, font_face="Calibri")
    add_text(s7, label, x_s+0.1, 2.15, 2.7, 0.85, 11, DARK_TEXT,
             align=PP_ALIGN.CENTER, font_face="Calibri")

# Left: Commercial + Government
add_rect(s7, 0.4, 3.3, 5.9, 3.9, WHITE)
add_rect(s7, 0.4, 3.3, 0.07, 3.9, GOLD)
add_text(s7, "[ 商業 ]  商業市場", 0.6, 3.45, 5.5, 0.4, 13, DARK_TEXT, bold=True, font_face="Calibri")
comm_items = [
    "美國商業收入 Q4 增長 137%",
    "客戶數量增長 34%",
    "IRS（美國稅務局）委託建置審計 AI 工具",
    "陸軍簽署五年期大合約，擴大 AI 平台使用範圍",
]
y_c = 3.95
for ci in comm_items:
    add_text(s7, f"▸  {ci}", 0.6, y_c, 5.5, 0.42, 12, DARK_TEXT, font_face="Calibri")
    y_c += 0.45

add_rect(s7, 0.4, 5.65, 5.9, 0.03, RGBColor(0xE2, 0xE8, 0xF0))
add_text(s7, "[ 政府 ]  軍政市場（護城河）", 0.6, 5.78, 5.5, 0.4, 13, DARK_TEXT, bold=True, font_face="Calibri")
gov_items = [
    "Maven 合約：2024年 $4.8億 → 2025年 $13億",
    "2026年3月正式列為全軍種標配項目記錄",
    "20年深耕五角大廈，頂級安全審計資質",
]
y_g = 6.28
for gi in gov_items:
    add_text(s7, f"▸  {gi}", 0.6, y_g, 5.5, 0.42, 12, DARK_TEXT, font_face="Calibri")
    y_g += 0.45

# Right: Moat analysis
add_rect(s7, 6.6, 3.3, 6.3, 3.9, DARK_BG)
add_rect(s7, 6.6, 3.3, 0.07, 3.9, TEAL)
add_text(s7, "[ 護城河 ]  Palantir 的護城河", 6.8, 3.45, 5.9, 0.4, 13, TEAL, bold=True, font_face="Calibri")

moat_items = [
    ("20 年安全審計資質", "每年通過審計、背景調查、合規檢查，非技術所能替代"),
    ("部署與人員訓練積累", "軍方人員已習慣 Palantir 介面與工作流程"),
    ("結構性政府綁定", "Maven 已列為全軍種強制配備，預算單獨列制"),
    ("Anthropic 自我設限", "Anthropic 與五角大廈關係惡化，短期難入政府市場"),
]
y_m = 3.95
for m_title, m_body in moat_items:
    add_text(s7, f"■  {m_title}", 6.8, y_m, 5.8, 0.35, 12, TEAL, bold=True, font_face="Calibri")
    add_text(s7, m_body, 6.8, y_m+0.35, 5.8, 0.42, 11, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_m += 0.9

# Footer
add_rect(s7, 0, 7.15, SLIDE_W, 0.35, NAVY)
add_text(s7, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s7, "7 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 世紀對決：蒂爾 vs 伯里
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
set_bg(s8, DARK_BG)

add_rect(s8, 0, 0, SLIDE_W, 1.1, CARD_BG)
add_rect(s8, 0, 1.1, SLIDE_W, 0.06, MINT)
add_text(s8, "07  /  世紀對決", 0.5, 0.1, 4, 0.35, 11, MINT, bold=True, font_face="Calibri")
add_text(s8, "彼得·蒂爾 VS Michael Burry：兩位傳奇的截然相反押注", 0.5, 0.38, 12, 0.6, 24, WHITE, bold=True, font_face="Calibri")

# VS divider
add_rect(s8, 6.45, 1.25, 0.43, 6.0, CARD_BG)
add_rect(s8, 6.53, 1.25, 0.07, 6.0, MINT)
add_text(s8, "VS", 6.35, 3.85, 0.83, 0.5, 16, MINT, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")

# Thiel side (BULL / LONG)
add_rect(s8, 0.3, 1.25, 6.0, 5.95, CARD_BG)
add_rect(s8, 0.3, 1.25, 0.07, 5.95, GOLD)

add_text(s8, "[ 成長 ]  彼得·蒂爾（看多）", 0.5, 1.38, 5.7, 0.42, 14, GOLD, bold=True, font_face="Calibri")
add_text(s8, "賭的是「位置」與「國家安全」", 0.5, 1.82, 5.7, 0.38, 12, RGBColor(0xCB, 0xD5, 0xE1),
         italic=True, font_face="Calibri")

add_rect(s8, 0.45, 2.3, 5.7, 0.03, rgb(0x30, 0x4A, 0x60))

thiel_items = [
    ("行動", "賣掉特斯拉、蘋果、微軟等所有持股\n唯獨死守 Palantir 股權一股未動"),
    ("核心邏輯", "模型可以換，位置換不了\nPalantir 靠的是一代人的信任資本與部署經驗"),
    ("護城河論點", "五角大廈安全許可資質需逐年審計累積\nAnthropic 短期內絕對無法複製"),
    ("供應商衝突利多", "Anthropic 與五角大廈關係惡化\n反而讓 Palantir 政府地位更加穩固"),
    ("目標", "軍方與政府市場「剛需 + 不可替代性」\n支撐長期高估值"),
]
y_t = 2.45
for title_t, body_t in thiel_items:
    add_text(s8, f"▶  {title_t}", 0.5, y_t, 5.6, 0.35, 11, GOLD, bold=True, font_face="Calibri")
    add_text(s8, body_t, 0.5, y_t+0.35, 5.6, 0.52, 11, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_t += 0.95

# Burry side (BEAR / SHORT)
add_rect(s8, 7.03, 1.25, 6.0, 5.95, CARD_BG)
add_rect(s8, 7.03, 1.25, 0.07, 5.95, RED_WARN)

add_text(s8, "📉  Michael Burry（做空）", 7.23, 1.38, 5.7, 0.42, 14, RED_WARN, bold=True, font_face="Calibri")
add_text(s8, "賭的是「競爭」與「估值重塑」", 7.23, 1.82, 5.7, 0.38, 12, RGBColor(0xCB, 0xD5, 0xE1),
         italic=True, font_face="Calibri")

add_rect(s8, 7.18, 2.3, 5.7, 0.03, rgb(0x30, 0x4A, 0x60))

burry_items = [
    ("行動", "持有大量 2027年到期、行使價 $50 美元的\n長期看跌期權（空頭倉位金額達數億美元）"),
    ("核心邏輯", "Anthropic 正在侵蝕 Palantir 的商業客戶\n剩餘政府業務撐不起 $3000億市值"),
    ("估值論點", "目前以「AI 平台獨角獸」定價（70倍 P/S）\n若降級為「國防 IT 供應商」：20倍或更低"),
    ("目標價", "$46 美元（較當時約 $130 低了近 2/3）"),
    ("RAMP 數據", "企業首次購 AI 工具，70% 選 Anthropic\nPalantir 的商業市場正被快速蠶食"),
]
y_b = 2.45
for title_b, body_b in burry_items:
    add_text(s8, f"▶  {title_b}", 7.23, y_b, 5.6, 0.35, 11, RED_WARN, bold=True, font_face="Calibri")
    add_text(s8, body_b, 7.23, y_b+0.35, 5.6, 0.52, 11, RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")
    y_b += 0.95

# Footer
add_rect(s8, 0, 7.15, SLIDE_W, 0.35, CARD_BG)
add_text(s8, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s8, "8 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 核心矛盾：估值定價之戰
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
set_bg(s9, LIGHT_BG)

add_rect(s9, 0, 0, SLIDE_W, 1.1, NAVY)
add_rect(s9, 0, 1.1, SLIDE_W, 0.06, GOLD)
add_text(s9, "08  /  核心矛盾", 0.5, 0.1, 4, 0.35, 11, GOLD, bold=True, font_face="Calibri")
add_text(s9, "3000億市值，給哪個故事定價？", 0.5, 0.38, 11, 0.6, 28, WHITE, bold=True, font_face="Calibri")

# Center question
add_rect(s9, 1.5, 1.35, 10.3, 0.85, DARK_BG)
add_text(s9,
    "「Palantir 目前的股價，是在為『AI 平台獨角獸』定價？還是為『國防剛需供應商』定價？」",
    1.65, 1.42, 10.0, 0.7, 14, GOLD, bold=False, italic=True,
    align=PP_ALIGN.CENTER, font_face="Calibri")

# Two scenario cards
# Scenario A: AI Platform
add_rect(s9, 0.4, 2.45, 6.0, 4.45, WHITE)
add_rect(s9, 0.4, 2.45, 0.08, 4.45, RED_WARN)
add_text(s9, "情境 A：若市場認定 AI 平台故事成立", 0.62, 2.58, 5.7, 0.38,
         13, RED_WARN, bold=True, font_face="Calibri")
add_text(s9, "👉 Burry 可能是對的", 0.62, 3.02, 5.7, 0.35, 12, RED_WARN, italic=True, font_face="Calibri")
a_lines = [
    "Anthropic 的 Managed Agents 已出現更便宜、更高效的替代方案",
    "商業市場客戶遷移成本低，可以今天用 Palantir、明天就換",
    "若商業收入被侵蝕，剩餘政府業務不支撐 70 倍 P/S 估值",
    "估值可能重塑至 20 倍甚至更低 → 目標價 $46",
    "Palantir 就像用米其林三星價格買了一個「雖好但不值那個價」的核桃",
]
y_a = 3.5
for al in a_lines:
    add_text(s9, f"▸  {al}", 0.62, y_a, 5.6, 0.5, 11, DARK_TEXT, font_face="Calibri")
    y_a += 0.53

# Scenario B: Defense Moat
add_rect(s9, 6.9, 2.45, 6.0, 4.45, WHITE)
add_rect(s9, 6.9, 2.45, 0.08, 4.45, GOLD)
add_text(s9, "情境 B：若市場認定國防剛需故事成立", 7.12, 2.58, 5.7, 0.38,
         13, GOLD, bold=True, font_face="Calibri")
add_text(s9, "👉 蒂爾可能是對的", 7.12, 3.02, 5.7, 0.35, 12, GOLD, italic=True, font_face="Calibri")
b_lines = [
    "五角大廈已將 Maven 列為全軍標配，預算結構性綁定",
    "政府合約替換成本極高：安全審計、人員訓練、部署習慣",
    "Anthropic 被列為供應鏈風險，短期難以進入政府市場",
    "20 年的信任積累是 Palantir 真正的護城河",
    "壟斷型市場 + 剛需屬性 → 高估值有其合理性",
]
y_b = 3.5
for bl in b_lines:
    add_text(s9, f"▸  {bl}", 7.12, y_b, 5.7, 0.5, 11, DARK_TEXT, font_face="Calibri")
    y_b += 0.53

# Bottom: author view
add_rect(s9, 0.4, 7.0, SLIDE_W-0.8, 0.4, DARK_BG)
add_text(s9,
    "影片作者觀點：兩者各有半對——商業市場難擋，政府市場難破。關鍵在於市場怎麼定義 Palantir 是「哪種公司」。",
    0.6, 7.03, 12.5, 0.35, 11, MINT, italic=True, font_face="Calibri")

# Footer
add_rect(s9, 0, 7.15, SLIDE_W, 0.35, NAVY)
add_text(s9, "Palantir (PLTR) 生存危機分析", 0.4, 7.18, 8, 0.28, 10, MUTED, font_face="Calibri")
add_text(s9, "9 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 結論：關鍵節點
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
set_bg(s10, DARK_BG)

# Left accent
add_rect(s10, 0, 0, 0.35, SLIDE_H, MINT)
# Top bar
add_rect(s10, 0.35, 0, SLIDE_W-0.35, 0.1, NAVY)
# Bottom bar
add_rect(s10, 0, 6.9, SLIDE_W, 0.6, NAVY)

add_text(s10, "09  /  結論", 0.6, 0.2, 3, 0.35, 11, MINT, bold=True, font_face="Calibri")
add_text(s10, "關鍵節點：5月11日財報見分曉", 0.6, 0.6, 10, 0.6, 30, WHITE, bold=True, font_face="Calibri")

add_rect(s10, 0.6, 1.3, 8.0, 0.04, TEAL)

# Three conclusion blocks
conclusions = [
    ("商業市場", "難以抵擋",
     "Anthropic Managed Agents 模式更便宜、更高效。\n"
     "對中型企業而言，開 API 帳戶遠比建 Palantir 平台更划算。\n"
     "商業客戶遷移門檻低，被蠶食的風險是真實存在的。",
     RED_WARN),
    ("政府/軍方市場", "護城河深厚",
     "20 年安全審計資質與部署習慣，Anthropic 短期複製不了。\n"
     "Maven 已列全軍標配，預算結構性綁定，不受短期競爭影響。\n"
     "Anthropic 自身與五角大廈關係緊張，反而鞏固了 Palantir 地位。",
     GOLD),
    ("股價核心問題", "估值定義之戰",
     "Palantir 的股價只有一個，但它同時扮演兩種角色。\n"
     "市場如何重新定義它，決定了蒂爾或 Burry 誰是對的。\n"
     "5月11日財報：觀察商業收入是否真的被 Anthropic 侵蝕。",
     MINT),
]
y_c = 1.5
for i, (area, verdict, body, accent) in enumerate(conclusions):
    x_c = 0.6 + i * 4.2
    w_c = 4.0
    add_rect(s10, x_c, y_c, w_c, 4.8, CARD_BG)
    add_rect(s10, x_c, y_c, w_c, 0.07, accent)
    add_rect(s10, x_c, y_c, 0.07, 4.8, accent)

    add_text(s10, area, x_c+0.15, y_c+0.12, w_c-0.25, 0.4, 14, accent, bold=True, font_face="Calibri")

    # Verdict badge
    add_rect(s10, x_c+0.15, y_c+0.58, w_c-0.3, 0.38, rgb(0x06, 0x5A, 0x82))
    add_text(s10, f"結論：{verdict}", x_c+0.18, y_c+0.6, w_c-0.35, 0.34,
             12, WHITE, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")

    add_text(s10, body, x_c+0.15, y_c+1.1, w_c-0.25, 3.5, 12,
             RGBColor(0xCB, 0xD5, 0xE1), font_face="Calibri")

# Countdown box
add_rect(s10, 0.6, 6.45, SLIDE_W-1.0, 0.38, rgb(0x06, 0x5A, 0x82))
add_text(s10,
    "[ 日程 ]  關注時間表：5月11日 PLTR 財報  |  蒂爾 13F 文件  |  Anthropic IPO 進展",
    0.8, 6.48, 12.0, 0.32, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, font_face="Calibri")

# Footer
add_rect(s10, 0, 7.15, SLIDE_W, 0.35, rgb(0x0A, 0x15, 0x20))
add_text(s10, "Palantir (PLTR) 生存危機分析  |  資料來源：YouTube 影片 / NotebookLM",
         0.4, 7.18, 10, 0.28, 10, MUTED, font_face="Calibri")
add_text(s10, "10 / 10", 12.5, 7.18, 0.7, 0.28, 10, MUTED, align=PP_ALIGN.RIGHT, font_face="Calibri")

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "/sessions/gracious-eager-dijkstra/mnt/outputs/PLTR_分析報告.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
